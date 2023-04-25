# 使用import导入requests模块
import requests
# TODO 使用import导入json模块
import json
# 使用from...import从pptx模块中导入Presentation
from pptx import Presentation
# TODO 使用import导入time模块
import time
import collections.abc

# 将.pptx文件路径赋值给变量path
path = "C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/5051/lesson.pptx"
# 读取path并赋值给变量pptxFile
pptxFile = Presentation(path)

# for循环遍历pptxFile中的.slides属性，赋值给slide
for slide in pptxFile.slides:
    # for循环遍历slide中.shapes属性，赋值给变量shape
    for shape in slide.shapes:
        # 判断形状中是否有文本框
        if shape.has_text_frame == True:
            # 读取形状中的文本框，并赋值给变量textFrame
            textFrame = shape.text_frame        
            # for循环遍历文本框内的所有段落
            # 赋值给变量paragraph
            for paragraph in textFrame.paragraphs:
                # for循环遍历段落中的所有样式块
                # 赋值给变量run
                for run in paragraph.runs:
                    # 读取样式块中的文本内容
                    texts = run.text
                    if texts != "":
                        # TODO 停顿1秒
                        time.sleep(1)
                        # TODO 构建一个字典word，添加键值为"i"，其值为texts;添加键值为"doctype"其值为"json"
                        word = {"i":texts,"doctype":"json"}
                        # 定义一个url,其值为"https://fanyi.youdao.com/translate"
                        url = "https://fanyi.youdao.com/translate"
                        # 将User-Agent以字典键对形式赋值给header
                        header = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_14_6) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/84.0.4147.89 Safari/537.36"}
                        # TODO 使用requests.post()，赋值给变量response
                        # 传入url
                        # 传入参数data，值为字典word
                        # 传如参数headers，值为字典header
                        response = requests.post(url,data=word,headers=header)
                        # TODO 使用.text获取内容并赋值给res_text
                        res_text = response.text
                        # TODO 使用json.loads()将str转换成dict，赋值给data
                        data = json.loads(res_text)
                        # 通过字典层层索引，找到翻译后的内容，并赋值给translate
                        translate = data["translateResult"][0][0]["tgt"]
                        # TODO 将该样式块的文本替换为翻译后的文本
                        run.text = translate
                        
 # TODO 将pptxFile保存至/yequ/new_lesson.pptx                     
pptxFile.save("C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/5051/new_lesson.pptx")