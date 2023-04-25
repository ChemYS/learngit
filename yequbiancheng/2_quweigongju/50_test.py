# PPT内英文翻译
#AttributeError: module 'collections' has no attribute 'Container'
#原因是python 3.10版本支持问题，此时在开头多导入一个依赖包collections.abc即可解决。
import collections.abc
# TODO 使用from...import从pptx模块中导入Presentation
from pptx import Presentation
# 将.pptx文件路径赋值给变量path
path = 'C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/5051/lesson.pptx'
# TODO 读取path并赋值给变量pptxFile
pptxFile = Presentation(path)
#print(pptxFile)
# TODO 导入requests模块
import requests
# TODO for循环遍历pptxFile中的.slides属性，赋值给slide
for slide in pptxFile.slides:
    # TODO for循环遍历slide中.shapes属性，赋值给变量shape
    #print(slide)
    for shape in slide.shapes:
        # TODO 判断形状中是否有文本框
        if shape.has_text_frame == True:
            # TODO 读取形状中的文本框，并赋值给变量textFrame
            textFrame = shape.text_frame
            # TODO for循环遍历文本框内的所有段落
            # 赋值给变量paragraph
            for paragraph in textFrame.paragraphs:
                # TODO for循环遍历段落中的所有样式块
                # 赋值给变量run
                for run in paragraph.runs:
                    # TODO 读取样式块中的文本内容
                    texts = run.text
                    if texts != "":
                        # TODO 构建一个字典word，添加键值为"i"其值为texts
                        # 添加键值为"doctype"其值为"json"。
                        word = {"i":texts,"doctype":"json"}
                        url = "https://fanyi.youdao.com/translate"
                        #url = "https://dict.youdao.com/webtranslate"
                        #将User-Agent以字典键对形式赋值给header
                        header = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36 Edg/112.0.1722.48"}
                        # TODO 使用requests.post()，赋值给变量response
                        # 传入url
                        # 传入参数data，值为字典word
                        # 传入参数headers，值为字典header
                        response = requests.post(url,data=word,headers=header)
                        # TODO 使用.text获取内容并赋值给res_text
                        res_text = response.text
                        print(res_text)