#冲哇！
import xlrd
import os
from pptx import Presentation
import collections.abc

path1 = "C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/39/candidate.xlsx"
table = xlrd.open_workbook(path1).sheets()[0]

path2 = "C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/39/picture"
allp = os.listdir(path2)

path3 = "C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/39/diploma.pptx"
pptxFile = Presentation(path3)
layoutOne = pptxFile.slide_layouts[0]

for row in range(1,5):
    e_name = table.row_values(row)[0]
    e_honor = table.row_values(row)[1]
    for i in allp:
        name = i.split(".")[0]
        if name == e_name:
            slide = pptxFile.slides.add_slide(layoutOne)
            #print(slide.placeholders)
            slide.placeholders[26].text = e_name
            slide.placeholders[0].text = e_honor
            slide.placeholders[27].insert_picture(path2 + "/" + i)

pptxFile.save("C:/Users/ChemYS/learngit/yequbiancheng/2_quweigongju/39/diploma_new.pptx")

'''
slide = pptxFile.slides
for single_plh in slide.placeholders:
    # 访问占位符对象中的.placeholder_format属性
    # 赋值给变量phf
    phf = single_plh.placeholder_format
    # TODO 访问phf中的.idx和.type属性
    # 格式化字符串以"-"连接并输出
    print(f"{phf.idx}-{phf.type}")
'''